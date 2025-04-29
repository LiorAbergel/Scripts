import json
import os
from pptx import Presentation
from deep_translator import GoogleTranslator

def extract_text_runs_to_json(pptx_path, json_path):
    """
    Extract text at the run level from each text frame in the PPTX and save
    the text along with its location (slide, shape, paragraph, and run index) to a JSON file.
    
    The JSON structure will be:
    {
        "slides": [
            {
                "slide_number": 1,
                "shapes": [
                    {
                        "shape_id": 123,
                        "paragraphs": [
                            {
                                "paragraph_index": 0,
                                "runs": [
                                    {
                                        "run_index": 0,
                                        "original_text": "Hello, ",
                                        "translated_text": "[Translated text here]"
                                    },
                                    {
                                        "run_index": 1,
                                        "original_text": "World!",
                                        "translated_text": "[Translated text here]"
                                    }
                                ]
                            },
                            ...
                        ]
                    },
                    ...
                ]
            },
            ...
        ]
    }
    
    Args:
        pptx_path (str): Path to the input PPTX file.
        json_path (str): Path to save the output JSON file.
    """
    prs = Presentation(pptx_path)
    data = {"slides": []}
    
    # Iterate over slides (slides are numbered starting at 1)
    for slide_index, slide in enumerate(prs.slides):
        slide_number = slide_index + 1
        slide_data = {"slide_number": slide_number, "shapes": []}
        
        # Loop through each shape in the slide
        for shape in slide.shapes:
            # Process only shapes with a text frame that contains text.
            if hasattr(shape, "has_text_frame") and shape.has_text_frame and shape.text_frame:
                shape_data = {"shape_id": shape.shape_id, "paragraphs": []}
                for p_index, paragraph in enumerate(shape.text_frame.paragraphs):
                    paragraph_data = {"paragraph_index": p_index, "runs": []}
                    for r_index, run in enumerate(paragraph.runs):
                        run_data = {
                            "run_index": r_index,
                            "original_text": run.text,
                            "translated_text": "[Translated text here]"
                        }
                        paragraph_data["runs"].append(run_data)
                    shape_data["paragraphs"].append(paragraph_data)
                slide_data["shapes"].append(shape_data)
        
        data["slides"].append(slide_data)
    
    # Ensure the output directory exists.
    os.makedirs(os.path.dirname(json_path), exist_ok=True)
    
    # Write the structured data to a JSON file.
    with open(json_path, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=4)
    
    print(f"Extraction complete. JSON saved to: {json_path}")


def insert_translated_text_runs(pptx_path, json_path, output_pptx_path):
    """
    Reads the JSON file (which now includes translated text for each run) and
    re-inserts the translated text into the corresponding locations in the PPTX,
    preserving the original run formatting, and ensuring all text is set to right-to-left.
    
    If a run's translated text is still the placeholder "[Translated text here]",
    the script will skip updating that run.
    
    Args:
        pptx_path (str): Path to the original PPTX file.
        json_path (str): Path to the JSON file with translations.
        output_pptx_path (str): Path to save the updated PPTX file.
    """
    # Load the JSON data with the translations.
    with open(json_path, "r", encoding="utf-8") as f:
        data = json.load(f)
        
    prs = Presentation(pptx_path)
    
    # Import the helper for modifying the underlying XML.
    from pptx.oxml.ns import qn

    def set_paragraph_rtl(paragraph):
        """
        Sets the paragraph to right-to-left by modifying its underlying XML.
        """
        pPr = paragraph._p.get_or_add_pPr()
        pPr.set(qn('a:rtl'), "1")
    
    # Iterate over each slide in the JSON data.
    for slide_data in data["slides"]:
        slide_number = slide_data["slide_number"]
        # Adjust index (Python list is 0-indexed while slide_number starts at 1)
        if slide_number - 1 >= len(prs.slides):
            print(f"Warning: Slide number {slide_number} not found in presentation.")
            continue
        slide = prs.slides[slide_number - 1]
        
        # Process each shape (text box) in the slide.
        for shape_data in slide_data["shapes"]:
            shape_id = shape_data["shape_id"]
            target_shape = None
            # Find the matching shape by shape_id.
            for shape in slide.shapes:
                if shape.shape_id == shape_id and hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    target_shape = shape
                    break
            if not target_shape:
                print(f"Warning: Shape with id {shape_id} not found on slide {slide_number}.")
                continue
            
            # Process paragraphs within the shape.
            for para_data in shape_data["paragraphs"]:
                para_index = para_data["paragraph_index"]
                if para_index >= len(target_shape.text_frame.paragraphs):
                    print(f"Warning: Paragraph index {para_index} not found in shape {shape_id} on slide {slide_number}.")
                    continue
                paragraph = target_shape.text_frame.paragraphs[para_index]
                
                # Process runs within the paragraph.
                for run_data in para_data["runs"]:
                    run_index = run_data["run_index"]
                    if run_index >= len(paragraph.runs):
                        print(f"Warning: Run index {run_index} not found in paragraph {para_index} in shape {shape_id} on slide {slide_number}.")
                        continue
                    translated_text = run_data.get("translated_text")
                    if translated_text is None:
                        print(f"Warning: Translated text is None for slide {slide_number}, shape {shape_id}, paragraph {para_index}, run {run_index}. Using original text instead.")
                        translated_text = run_data.get("original_text", "")
                    # Skip updating this run if the translated text is still the placeholder.
                    if translated_text.strip() == "[Translated text here]":
                        continue
                    # Update the run's text with the translated text.
                    paragraph.runs[run_index].text = translated_text
                
                # Set the paragraph to right-to-left.
                set_paragraph_rtl(paragraph)
                    
    prs.save(output_pptx_path)
    print(f"Insertion complete. Updated presentation saved to: {output_pptx_path}")


def auto_translate_json(json_path, output_json_path, target_language="hebrew"):
    """
    Reads the JSON file with the original text and automatically translates
    each text run using GoogleTranslator from deep_translator.
    
    The translated text is saved in the 'translated_text' field for each run.
    
    Args:
        json_path (str): Path to the JSON file with original text.
        output_json_path (str): Path to save the updated JSON file with translations.
        target_language (str): The target language (e.g., 'hebrew').
    """
    # Load the JSON data.
    with open(json_path, "r", encoding="utf-8") as f:
        data = json.load(f)
    
    # Initialize the translator.
    translator = GoogleTranslator(source='auto', target=target_language)
    
    # Iterate over each text run and translate.
    for slide in data["slides"]:
        for shape in slide["shapes"]:
            for paragraph in shape["paragraphs"]:
                for run in paragraph["runs"]:
                    original_text = run["original_text"]
                    if original_text.strip():  # Only translate non-empty text
                        try:
                            translated = translator.translate(text=original_text)
                            run["translated_text"] = translated
                        except Exception as e:
                            print(f"Error translating text: {original_text}\nError: {e}")
                            # Fallback: keep the original text
                            run["translated_text"] = original_text
    
    # Ensure the output directory exists.
    os.makedirs(os.path.dirname(output_json_path), exist_ok=True)
    
    # Save the updated JSON data.
    with open(output_json_path, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=4)
    
    print(f"Translation complete. Updated JSON saved to: {output_json_path}")


def main():
    """
    Main driver that allows you to extract text for translation, automatically
    translate the extracted text, or reinsert translated text into a PPTX file.
    """
    mode = input("Enter mode (extract/translate/insert): ").strip().lower()
    
    if mode == "extract":
        pptx_path = input("Enter the path to the PPTX file: ").strip('"')
        json_path = input("Enter the path to save the JSON file (e.g., translations.json): ").strip('"')
        if not os.path.exists(pptx_path):
            print(f"Error: File '{pptx_path}' does not exist.")
            return
        extract_text_runs_to_json(pptx_path, json_path)
    
    elif mode == "translate":
        json_input = input("Enter the path to the JSON file with original text: ").strip('"')
        output_json = input("Enter the path to save the translated JSON file: ").strip('"')
        target_language = input("Enter the target language (e.g., 'hebrew'): ").strip() or "hebrew"
        if not os.path.exists(json_input):
            print(f"Error: JSON file '{json_input}' does not exist.")
            return
        auto_translate_json(json_input, output_json, target_language)
    
    elif mode == "insert":
        pptx_path = input("Enter the path to the original PPTX file: ").strip('"')
        json_path = input("Enter the path to the JSON file with translations: ").strip('"')
        output_pptx_path = input("Enter the path to save the updated PPTX file: ").strip('"')
        if not os.path.exists(pptx_path):
            print(f"Error: File '{pptx_path}' does not exist.")
            return
        if not os.path.exists(json_path):
            print(f"Error: File '{json_path}' does not exist.")
            return
        insert_translated_text_runs(pptx_path, json_path, output_pptx_path)
    
    else:
        print("Invalid mode. Please enter 'extract', 'translate', or 'insert'.")


if __name__ == "__main__":
    main()
